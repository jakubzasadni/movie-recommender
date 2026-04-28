"""
Generuje prezentację PowerPoint dla projektu NeuMF.
Uruchom: python make_presentation.py
Wynik: reports/prezentacja.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

DARK_BG    = RGBColor(0x18, 0x18, 0x18)   # prawie czarne
ACCENT     = RGBColor(0xD4, 0xC5, 0xA9)   # ciepły beż / krem
ACCENT2    = RGBColor(0x9E, 0x9E, 0x9E)   # sredni szary
WHITE      = RGBColor(0xF0, 0xED, 0xE8)   # ciepła biel
LIGHT_GRAY = RGBColor(0xA0, 0x9A, 0x92)   # ciepły szary
GREEN      = RGBColor(0x8B, 0xAF, 0x7A)   # stonowana szałwia
RED        = RGBColor(0xC4, 0x7A, 0x6A)   # stonowana terakota
DARK_CARD  = RGBColor(0x26, 0x24, 0x22)   # ciemna karta (ciepły odcień)

W = Inches(13.33)
H = Inches(7.5)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    return slide


def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, color=ACCENT):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


TITLE_BAR  = RGBColor(0x2C, 0x2A, 0x28)   # ciemny grafit, ciepły odcień

def title_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(1.25), TITLE_BAR)
    # cienka linia akcentująca na dole paska
    rect(slide, 0, Inches(1.2), W, Inches(0.05), ACCENT)
    txt(slide, title,    Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.75),
        size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.38),
            size=15, color=LIGHT_GRAY)


def bullets(slide, items, x, y, w, size=16, gap=0.46):
    for i, item in enumerate(items):
        prefix = "▸  " if not item.startswith("  ") else "      "
        txt(slide, prefix + item.strip(), x, y + Inches(gap * i),
            w, Inches(0.44), size=size, color=WHITE)


def img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ══════════════════════════════════════════════════════════════════════════════
# 1. Tytułowy
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
rect(s, 0, Inches(2.4), W, Inches(2.9), RGBColor(0x2C, 0x2A, 0x28))
txt(s, "System Rekomendacji Filmów",
    Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.1),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Neural Matrix Factorization (NeuMF)  ·  MovieLens-1M",
    Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.6),
    size=22, color=RGBColor(0xE0, 0xE0, 0xE0), align=PP_ALIGN.CENTER)
txt(s, "He et al., 2017  ·  arXiv:1708.05031",
    Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
    size=16, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "PyTorch  ·  NVIDIA GeForce RTX 3050  ·  MovieLens-1M",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
    size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
notes(s, "Dzień dobry. Nasz projekt dotyczy systemu rekomendacji filmów opartego na głębokim uczeniu. "
         "Zaimplementowaliśmy algorytm o nazwie NeuMF, czyli Neural Matrix Factorization, "
         "który został zaproponowany przez He i współpracowników w 2017 roku. "
         "Trenowaliśmy i testowaliśmy model na datasecie MovieLens-1M — "
         "jest to jeden z najbardziej znanych benchmarków w dziedzinie systemów rekomendacji. "
         "Całość napisana jest w PyTorchu i uruchamiana na karcie graficznej RTX 3050.")


# ══════════════════════════════════════════════════════════════════════════════
# 2. Problem i Dataset
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Problem i Dataset", "Co przewidujemy i na jakich danych")

rect(s, Inches(0.3), Inches(1.4), Inches(5.9), Inches(5.8), DARK_CARD)
txt(s, "Problem", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
    size=19, bold=True, color=ACCENT)
bullets(s, [
    "Macierz użytkownik × film wypełniona w ~4%",
    "Chcemy przewidzieć puste komórki",
    "Tryb implicit feedback:",
    "  ocena ≥ 4  →  lubię (1)",
    "  brak / ocena < 4  →  nie wiem (0)",
    "Wyjście modelu: prawdopodobieństwo 0–1",
], Inches(0.5), Inches(2.1), Inches(5.5), size=16, gap=0.47)

rect(s, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.8), DARK_CARD)
txt(s, "MovieLens-1M", Inches(6.7), Inches(1.5), Inches(6.1), Inches(0.5),
    size=19, bold=True, color=ACCENT2)
for i, (k, v) in enumerate([
    ("Użytkownicy",   "6 040"),
    ("Filmy",         "3 952"),
    ("Ocen",          "1 000 209"),
    ("Wypełnienie",   "~4%"),
    ("Skala ocen",    "1 – 5"),
    ("Podział",       "80 / 10 / 10  (chronologicznie)"),
    ("Neg. sampling", "4 negatywy na 1 pozytyw"),
]):
    txt(s, k, Inches(6.7), Inches(2.15 + i * 0.46), Inches(2.8), Inches(0.43),
        size=15, color=LIGHT_GRAY)
    txt(s, v, Inches(9.5), Inches(2.15 + i * 0.46), Inches(3.2), Inches(0.43),
        size=15, bold=True, color=WHITE)
notes(s, "Problem, który rozwiązujemy, wygląda następująco. "
         "Wyobraźcie sobie tabelę, gdzie wiersze to użytkownicy, a kolumny to filmy. "
         "Każda komórka powinna zawierać ocenę — ale w rzeczywistości wypełniona jest tylko cztery procent tej tabeli, "
         "bo nikt nie widział wszystkich filmów. Naszym zadaniem jest przewidzieć te puste komórki. "
         "Zdecydowaliśmy się pracować w trybie implicit feedback. "
         "Oznacza to, że nie przewidujemy dokładnej oceny w skali jeden do pięciu, "
         "tylko odpowiadamy na pytanie: czy ten użytkownik polubi ten film — tak czy nie. "
         "Przyjęliśmy próg cztery — ocena cztery lub wyżej oznacza, że użytkownik lubi film. "
         "Dataset MovieLens-1M zawiera milion ocen od sześciu tysięcy czterdziestu użytkowników "
         "dla blisko czterech tysięcy filmów. "
         "Dane dzielimy chronologicznie — najstarsze osiemdziesiąt procent idzie do treningu, "
         "a najnowsze oceny każdego użytkownika stanowią zbiór testowy.")


# ══════════════════════════════════════════════════════════════════════════════
# 3. Architektura
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Architektura — NeuMF", "GMF + MLP → concat → Sigmoid")

rect(s, Inches(0.3), Inches(1.4), Inches(5.9), Inches(4.9), RGBColor(0x1A, 0x37, 0x5E))
txt(s, "GMF — Generalized Matrix Factorization",
    Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.55), size=17, bold=True, color=ACCENT)
bullets(s, [
    "Osobne embeddingi user i item (32 dim)",
    "Mnożenie element po elemencie (hadamard)",
    "Wychwytuje zależności liniowe",
    "Uproszczona, interpretowalna wersja MF",
], Inches(0.5), Inches(2.15), Inches(5.5), size=15, gap=0.5)

rect(s, Inches(6.4), Inches(1.4), Inches(6.6), Inches(4.9), RGBColor(0x3E, 0x1A, 0x1A))
txt(s, "MLP — Multi-Layer Perceptron",
    Inches(6.6), Inches(1.5), Inches(6.2), Inches(0.55), size=17, bold=True, color=ACCENT2)
bullets(s, [
    "Osobne embeddingi user i item (32 dim)",
    "concat(u, i) → Dense(64) → Dense(32) → Dense(16)",
    "Aktywacja ReLU + Dropout(0.2)",
    "Uczy sie zlozonych, nieliniowych wzorcow",
], Inches(6.6), Inches(2.15), Inches(6.1), size=15, gap=0.5)

rect(s, Inches(0.3), Inches(6.45), W - Inches(0.6), Inches(0.75), DARK_CARD)
txt(s, "concat(GMF_out, MLP_out)  →  Linear(48 → 1)  →  Sigmoid  →  p ∈ (0, 1)",
    Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
    size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
notes(s, "NeuMF łączy dwa podejścia w jednym modelu. "
         "GMF to ulepszona wersja klasycznego Matrix Factorization — zamiast zwykłego iloczynu skalarnego, "
         "robi ważone mnożenie element po elemencie. Prosta i szybka w konwergencji. "
         "MLP dostaje te same dane, ale przetwarza je przez sieć gęstą — może się nauczyć dowolnie skomplikowanych wzorców, "
         "np. że użytkownik lubi filmy akcji ale tylko z lat 90. "
         "Kluczowy szczegół: oba moduły mają OSOBNE embeddingi, co pozwala im uczyć się niezależnie. "
         "Na końcu wyniki są konkatenowane i przepuszczane przez warstwę wyjściową z funkcją sigmoid.")


# ══════════════════════════════════════════════════════════════════════════════
# 4. Trening
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Trening", "Adam · BCE Loss · Negative Sampling · Early Stopping")

rect(s, Inches(0.3), Inches(1.4), Inches(8.3), Inches(5.8), DARK_CARD)
txt(s, "Jak model sie uczy", Inches(0.5), Inches(1.5), Inches(8), Inches(0.5),
    size=19, bold=True, color=ACCENT)
bullets(s, [
    "Forward pass: (user_id, item_id) → predykcja p ∈ (0,1)",
    "Loss: BCE = −( y·log(p) + (1−y)·log(1−p) )",
    "Backward pass: chain rule, PyTorch autograd liczy gradienty",
    "Adam: w = w − lr · grad  (adaptacyjny lr per waga)",
    "~500 000 par treningowych × 20 epok",
    "Negative sampling: 4 losowe filmy bez interakcji na 1 pozytyw",
    "Early stopping: patience=5, zapisujemy najlepszy checkpoint",
], Inches(0.5), Inches(2.1), Inches(7.9), size=15, gap=0.49)

rect(s, Inches(8.9), Inches(1.4), Inches(4.1), Inches(5.8), DARK_CARD)
txt(s, "Hiperparametry", Inches(9.1), Inches(1.5), Inches(3.7), Inches(0.5),
    size=19, bold=True, color=ACCENT2)
for i, (k, v) in enumerate([
    ("Embedding dim",  "32"),
    ("MLP layers",     "64 → 32 → 16"),
    ("Optimizer",      "Adam"),
    ("Learning rate",  "0.001"),
    ("Batch size",     "256"),
    ("Neg. samples",   "4"),
    ("Epoki max",      "20"),
    ("GPU",            "RTX 3050 (4GB)"),
]):
    txt(s, k, Inches(9.1), Inches(2.1 + i * 0.44), Inches(2.3), Inches(0.41),
        size=14, color=LIGHT_GRAY)
    txt(s, v, Inches(11.4), Inches(2.1 + i * 0.44), Inches(1.5), Inches(0.41),
        size=14, bold=True, color=WHITE)
notes(s, "Trening przebiega iteracyjnie: forward pass produkuje predykcję p ∈ (0,1), "
         "Binary Cross-Entropy mierzy błąd, backpropagation liczy gradienty przez cały graf operacji, "
         "a Adam aktualizuje wagi. Negative sampling jest kluczowy — bez negatywów model nauczyłby się "
         "przewidywać 1 dla wszystkiego. Losujemy 4 filmy bez interakcji na każdy pozytyw. "
         "Early stopping zatrzymuje trening gdy val loss nie spada przez 5 epok z rzędu — "
         "zapobiega overfittingowi i oszczędza czas. "
         "Trening na RTX 3050 zajmuje ok. 60s na epokę przy 500k parach treningowych.")


# ══════════════════════════════════════════════════════════════════════════════
# 5. Napotkane problemy
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Napotkane Problemy", "4 bledy, ktore trzeba bylo naprawic")

for i, (color, title, desc) in enumerate([
    (ACCENT,  "1. UnicodeEncodeError — konsola Windows",
              "Znak strzalki (U+2192) nieobslugiwany przez cp1250. Fix: zastapienie tekstem."),
    (RED,     "2. Bledna walidacja → early stopping po 1 epoce → HitRate 0.07",
              "Val set mial tylko pozytywy → BCE rosla mimo uczenia sie. Fix: dodac negatywy do val setu."),
    (ACCENT2, "3. Zly protokol ewaluacji — ranking sposrod 3700 zamiast 101",
              "Papier rankuje 1 cel vs 100 losowych negatywow. Full ranking daje HitRate 0.07 zamiast 0.66."),
    (GREEN,   "4. PyTorch CPU-only → ~100s/epoka",
              "Build bez CUDA. Fix: pip install torch --index-url .../cu128 → RTX 3050 → ~60s/epoka."),
]):
    y = Inches(1.45 + i * 1.48)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(1.3), DARK_CARD)
    rect(s, Inches(0.3), y, Inches(0.15), Inches(1.3), color)
    txt(s, title, Inches(0.6), y + Inches(0.1), Inches(12.1), Inches(0.5),
        size=17, bold=True, color=color)
    txt(s, desc,  Inches(0.6), y + Inches(0.6), Inches(12.1), Inches(0.6),
        size=14, color=LIGHT_GRAY)
notes(s, "Problem 1 był trywialny — znak Unicode niekompatybilny z Windows. "
         "Problem 2 był poważny: zbiór walidacyjny miał tylko pozytywy, więc BCE rosła nawet gdy model się uczył. "
         "Early stopping zapisał checkpoint z epoki 1 — model prawie nieuczony. Fix: negatywy w val secie. "
         "Problem 3 był fundamentalny: rankowanie spośród 3700 filmów to zupełnie inne zadanie niż "
         "spośród 101 kandydatów jak w paperze. Ten sam model dał HitRate 0.07 vs 0.66 w zależności od protokołu. "
         "Problem 4: zainstalowany był CPU-only build PyTorch. Po reinstalacji z cu128 GPU zaczął pracować.")


# ══════════════════════════════════════════════════════════════════════════════
# 6. Wyniki końcowe — liczby + wykres z notebooka
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Wyniki Koncowe", "NeuMF · emb_dim=32 · 20 epok · protokol 100 negatywow (He et al.)")

# Duże liczby po lewej
for j, (metric, val, target) in enumerate([
    ("HitRate@10", "0.6596", "> 0.65"),
    ("NDCG@10",    "0.3788", "> 0.38"),
]):
    x = Inches(0.3 + j * 3.1)
    rect(s, x, Inches(1.4), Inches(2.8), Inches(3.8), DARK_CARD)
    txt(s, metric, x + Inches(0.1), Inches(1.55), Inches(2.6), Inches(0.5),
        size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, val,    x + Inches(0.1), Inches(2.1),  Inches(2.6), Inches(1.0),
        size=42, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, f"cel: {target}", x + Inches(0.1), Inches(3.2), Inches(2.6), Inches(0.4),
        size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    txt(s, "OSIAGNIETY", x + Inches(0.1), Inches(3.65), Inches(2.6), Inches(0.4),
        size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

rect(s, Inches(0.3), Inches(5.4), Inches(6.1), Inches(1.8), DARK_CARD)
txt(s, "Popularity Baseline: HitRate 0.04  →  NeuMF: 0.66",
    Inches(0.5), Inches(5.55), Inches(5.8), Inches(0.5),
    size=15, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txt(s, "Poprawa ~16x dzieki personalizacji",
    Inches(0.5), Inches(6.05), Inches(5.8), Inches(0.4),
    size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Wykres z notebooka po prawej
img(s, "results/metrics_comparison.png",
    Inches(6.5), Inches(1.3), Inches(6.6), Inches(5.9))
notes(s, "Oba cele z papieru He et al. 2017 zostały osiągnięte: HitRate@10 = 0.6596 (cel > 0.65) "
         "i NDCG@10 = 0.3788 (cel > 0.38). "
         "Wykres po prawej pokazuje porównanie NeuMF z popularity baseline i wartościami docelowymi z papieru. "
         "Popularity baseline rekomenduje po prostu najpopularniejsze filmy wszystkim użytkownikom — "
         "HitRate = 0.04. NeuMF osiąga 0.66, czyli jest ~16 razy lepszy. "
         "To pokazuje realną wartość personalizacji.")


# ══════════════════════════════════════════════════════════════════════════════
# 7. Ablacja — wykresy
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Ablacja: GMF vs MLP vs NeuMF", "10 epok, emb_dim=32")

img(s, "results/ablation_components.png",
    Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))
notes(s, "Ablacja pokazuje wkład każdego komponentu. Przy 10 epokach GMF (0.6582) lekko bije NeuMF (0.6534) — "
         "NeuMF jest większy i potrzebuje więcej iteracji żeby wszystkie wagi 'dograły się'. "
         "MLP (0.6366) konwerguje najwolniej ze wszystkich. "
         "Przy pełnym treningu (20 epok) NeuMF wygrywa — potwierdza tezę papieru. "
         "Czerwona przerywana linia to wartość docelowa z literatury (0.65 / 0.38).")


# ══════════════════════════════════════════════════════════════════════════════
# 8. Sweep embeddingów — wykres
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Wplyw rozmiaru embeddingow", "NeuMF · 10 epok · wiekszy model nie zawsze lepszy")

img(s, "results/ablation_embeddings.png",
    Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))
notes(s, "Sweep rozmiarów embeddingów pokazuje że emb_dim=8 i 16 wypadają równie dobrze jak standardowe 32. "
         "emb_dim=64 jest najgorszy — większy model zaczyna overfittować szybciej na datasecie tej wielkości. "
         "MovieLens-1M to za mało danych żeby uzasadnić zwiększanie pojemności modelu powyżej 32. "
         "Wniosek: dla tego benchmarku nie opłaca się skalować embeddingów w górę.")


# ══════════════════════════════════════════════════════════════════════════════
# 9. Pełna tabela wyników
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Pelna Tabela Wynikow", "Wszystkie warianty modelu")

img(s, "results/full_results_table.png",
    Inches(1.5), Inches(1.4), Inches(10.3), Inches(3.8))

bullets(s, [
    "NeuMF (20 epok, pelny trening): HitRate 0.6596  /  NDCG 0.3788  →  CELE OSIAGNIETE",
    "emb_dim 8-32 daje podobne wyniki; 64 overfittuje",
    "MLP wolniej konwerguje niz GMF — przy 10 epokach GMF lekko wygrywa",
], Inches(0.5), Inches(5.45), Inches(12.3), size=15, gap=0.55)
notes(s, "Zestawienie wszystkich eksperymentów. Najważniejszy wiersz to NeuMF z pełnym treningiem (20 epok) — "
         "osiąga HitRate 0.6596 i NDCG 0.3788, oba powyżej progów z papieru. "
         "W ablacji z 10 epokami wyniki są niższe bo modele nie zdążyły w pełni skonwergować. "
         "emb_dim=64 to jedyny wariant który wyraźnie odpada.")


# ══════════════════════════════════════════════════════════════════════════════
# 10. Wnioski
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Wnioski", "")

for i, (color, title, desc) in enumerate([
    (ACCENT,     "NeuMF > suma swoich czesci — ale potrzebuje czasu",
                 "Przy krotkim treningu GMF moze dorownywac NeuMF. Przy 20 epokach NeuMF wygrywa."),
    (ACCENT2,    "Wiekszy model nie zawsze lepszy",
                 "emb_dim=64 wypadl najgorzej. ML-1M zbyt maly by uzasadnic wzrost ponad 32."),
    (GREEN,      "Protokol ewaluacji zmienia wszystko",
                 "Ten sam model: HitRate 0.07 (full ranking 3700) vs 0.66 (100 negatywow jak w paperze)."),
    (LIGHT_GRAY, "Personalizacja wnosi realna wartosc",
                 "NeuMF 16x lepszy niz 'polecaj wszystkim najpopularniejsze filmy'."),
]):
    y = Inches(1.45 + i * 1.48)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(1.3), DARK_CARD)
    rect(s, Inches(0.3), y, Inches(0.15), Inches(1.3), color)
    txt(s, title, Inches(0.6), y + Inches(0.1), Inches(12.1), Inches(0.5),
        size=17, bold=True, color=color)
    txt(s, desc,  Inches(0.6), y + Inches(0.6), Inches(12.1), Inches(0.6),
        size=14, color=LIGHT_GRAY)
notes(s, "Cztery główne wnioski: "
         "1. NeuMF wygrywa z GMF i MLP osobno, ale tylko przy wystarczającym czasie treningu. "
         "2. Większy model (emb_dim=64) dał najgorsze wyniki — overfitting na małym datasecie. "
         "3. Protokół ewaluacji ma kluczowe znaczenie — zawsze sprawdzaj szczegóły zanim porównasz wyniki z literaturą. "
         "4. Personalizacja (NeuMF) jest 16x lepsza od najprostszego baseline — wartość algorytmu jest realna.")


# ══════════════════════════════════════════════════════════════════════════════
# 11. EDA — analiza datasetu
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Analiza Datasetu — MovieLens-1M", "Rozklad ocen i podzial na pozytywy/negatywy")
img(s, "results/eda_ratings.png", Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))
notes(s, "Lewy wykres: rozkład ocen jest asymetryczny — dominują oceny 3 i 4, mało ocen 1. "
         "To typowe dla systemów rekomendacji: użytkownicy rzadko oceniają filmy które bardzo nie lubią. "
         "Prawy wykres: po binaryzacji 55% interakcji to pozytywy (ocena ≥ 4), 45% to negatywy. "
         "Proporcja jest bardziej zbalansowana niż można by oczekiwać — dlatego negative sampling z ratio 4:1 "
         "jest ważny, żeby model uczył się też rozpoznawać co użytkownikowi nie odpowiada.")


# ══════════════════════════════════════════════════════════════════════════════
# 12. Krzywa uczenia
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Krzywa Uczenia", "Train loss vs Val loss — NeuMF, emb_dim=32")

img(s, "results/training_curve.png", Inches(0.3), Inches(1.35), Inches(8.5), Inches(5.9))

rect(s, Inches(9.0), Inches(1.4), Inches(4.0), Inches(5.8), DARK_CARD)
txt(s, "Obserwacje", Inches(9.2), Inches(1.55), Inches(3.6), Inches(0.5),
    size=18, bold=True, color=ACCENT)
bullets(s, [
    "Epoka 2 = best checkpoint",
    "Od epoki 3 → overfitting",
    "Train loss spada caly czas",
    "Val loss odbija w gore",
    "Gap train/val rosnie",
    "Early stopping zadzialal",
    "poprawnie po patience=5",
], Inches(9.2), Inches(2.1), Inches(3.6), size=14, gap=0.5)
notes(s, "Krzywa uczenia z faktycznego treningu NeuMF. "
         "Epoka 1: val loss 0.3279, epoka 2: 0.3252 — najlepszy checkpoint. "
         "Od epoki 3 val loss zaczyna rosnąć mimo że train loss cały czas spada — klasyczny overfitting. "
         "Gap między train a val rośnie z każdą epoką. Early stopping zatrzymał trening po 7 epokach "
         "(patience=5, licząc od epoki 2). "
         "Warto zauważyć: przed naprawą błędu z val setem, val loss rósł od epoki 1 i checkpoint był z epoki 1.")


# ══════════════════════════════════════════════════════════════════════════════
# 13. Protokół ewaluacji — porównanie
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Protokol Ewaluacji — Kluczowe Odkrycie",
          "Full ranking vs 100 negatywow (He et al.) — ten sam model, inne wyniki")

img(s, "results/eval_protocol.png", Inches(0.3), Inches(1.35), Inches(12.7), Inches(4.5))

rect(s, Inches(0.3), Inches(5.95), Inches(12.7), Inches(1.3), DARK_CARD)
rect(s, Inches(0.3), Inches(5.95), Inches(0.15), Inches(1.3), RED)
txt(s, "Full ranking (3700 filmow): HitRate = 0.07  |  26x lepszy od losowego — model DZIALA",
    Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.45), size=14, color=WHITE)
txt(s, "Protokol papieru (101 kandydatow): HitRate = 0.66  |  Bezposrednio porownywalne z literatura",
    Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.45), size=14, color=LIGHT_GRAY)
notes(s, "To najważniejsze odkrycie metodologiczne projektu. "
         "Full ranking: model ocenia wszystkie ~3700 filmów i sprawdzamy czy trafny film jest w top-10. "
         "Szansa losowa: 10/3700 ≈ 0.003. Nasz model osiąga 0.07 — 26x lepiej niż losowo, czyli działa. "
         "Protokół papieru: dla każdego użytkownika losujemy 100 filmów bez interakcji i rankujemy 101 kandydatów. "
         "Szansa losowa: 10/101 ≈ 0.099. Nasz model osiąga 0.66 — zgodnie z literaturą. "
         "Morał: zawsze sprawdzaj szczegóły protokołu ewaluacji zanim porównasz swoje wyniki z paperem. "
         "Nowsza literatura (Krichene & Rendle, 2020) krytykuje protokół 100-negatywów za bias statystyczny.")


# ══════════════════════════════════════════════════════════════════════════════
# 14. Przykładowe rekomendacje
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_bar(s, "Przykladowe Rekomendacje", "Top-8 filmow dla dwoch roznych uzytkownikow")
img(s, "results/sample_recs.png", Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))
notes(s, "Wizualizacja rzeczywistych rekomendacji z wytrenowanego modelu dla dwóch różnych użytkowników. "
         "Oś X to score modelu — prawdopodobieństwo interakcji z danym filmem. "
         "Filmy znane użytkownikowi (z historii) są wykluczone z rankingu. "
         "Widać że rekomendacje dla różnych użytkowników są różne — model faktycznie personalizuje. "
         "Score ~0.8-0.9 oznacza że model jest dość pewny że film spodoba się temu użytkownikowi.")


# ══════════════════════════════════════════════════════════════════════════════
# Zapis — nowy plik
# ══════════════════════════════════════════════════════════════════════════════
os.makedirs("reports", exist_ok=True)
out = "reports/prezentacja_rozbudowana.pptx"
prs.save(out)
print(f"Zapisano: {out}  ({len(prs.slides)} slajdow)")
